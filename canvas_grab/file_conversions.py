import pypdfium2 as pdfium
import json
from bs4 import BeautifulSoup
from termcolor import colored
from docx import Document
from pptx import Presentation

def convert_file_to_json(path):
    if path.endswith('.pdf'):
        return pdf_to_json(path)
    elif path.endswith('.html'):
        return html_to_json(path)
    elif path.endswith('.docx'):
        return docx_to_json(path)
    elif path.endswith('.pptx'):
        return pptx_to_json(path)
    else:
        # Log a warning, use color
        print(colored(f'  Unsupported file type: {path}', 'yellow'))
        return None 

def pdf_to_json(path):
    pdf = pdfium.PdfDocument(path)
    document = {"document_name": path.split('/')[-1], "content": []}
    for i, page in enumerate(pdf, start=1):
        document["content"].append({"page_number": i, "text": page.get_textpage().get_text_bounded()})
    return json.dumps(document, indent=4)

def html_to_json(path):
    with open(path, 'r', encoding='utf-8') as file:
        html_content = file.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    document = {"document_name": path.split('/')[-1], "content": {}}
    
    # Check if the page contains a meta refresh (redirect)
    meta_refresh = soup.find('meta', attrs={"http-equiv": "refresh"})
    if meta_refresh:
        document["content"] = {
            "type": "redirect",
            "redirect_url": meta_refresh["content"].split("URL=")[-1]
        }
    else:
        # Extract the title, headings, and paragraphs if the page contains actual content
        document["content"] = {
            "title": soup.title.string if soup.title else None,
            "paragraphs": [p.get_text(strip=True) for p in soup.find_all('p')]
        }
    
    return json.dumps(document, indent=4)

def docx_to_json(path):
    # Load the .docx file
    doc = Document(path)
    document = {"document_name": path.split('/')[-1], "content": []}
    
    # Iterate over each paragraph and add it to content
    for i, paragraph in enumerate(doc.paragraphs, start=1):
        # Check if the paragraph has text (skip empty paragraphs)
        if paragraph.text.strip():
            document["content"].append({"paragraph_number": i, "text": paragraph.text.strip()})
    
    return json.dumps(document, indent=4)

def pptx_to_json(path):
    prs = Presentation(path)
    document = {"document_name": path.split('/')[-1], "content": []}
    
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Each shape might contain multiple paragraphs
                # You can either join them or store them as a list
                slide_text = ' '.join([p.text for p in shape.text_frame.paragraphs if p.text.strip()])
                if slide_text.strip():
                    slide_texts.append(slide_text.strip())
        
        # Only add entry if there is text on the slide
        if slide_texts:
            slide_texts = "\n".join(slide_texts)
            document["content"].append({"slide_number": i, "text": slide_texts})
    
    return json.dumps(document, indent=4)
